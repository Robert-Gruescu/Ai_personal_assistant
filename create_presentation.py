from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Alias for convenience
RgbColor = RGBColor

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(63, 81, 181)  # Indigo
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(200, 200, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=RgbColor(63, 81, 181)):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = accent_color
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet.startswith("##"):
            p.text = bullet[2:].strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = accent_color
            p.space_before = Pt(16)
        elif bullet.startswith("-"):
            p.text = "    • " + bullet[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = RgbColor(50, 50, 50)
            p.space_before = Pt(6)
        else:
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RgbColor(50, 50, 50)
            p.space_before = Pt(10)
    
    return slide

def add_architecture_slide(prs, title, components):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(63, 81, 181)
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Architecture boxes
    colors = [
        RgbColor(76, 175, 80),   # Green
        RgbColor(33, 150, 243),  # Blue  
        RgbColor(255, 152, 0),   # Orange
        RgbColor(156, 39, 176),  # Purple
        RgbColor(244, 67, 54),   # Red
        RgbColor(0, 188, 212),   # Cyan
    ]
    
    start_y = 1.6
    for i, (comp_title, comp_items) in enumerate(components):
        color = colors[i % len(colors)]
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3 + (i % 3) * 3.2),
            Inches(start_y + (i // 3) * 2.4),
            Inches(3),
            Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Box title
        text_box = slide.shapes.add_textbox(
            Inches(0.4 + (i % 3) * 3.2),
            Inches(start_y + 0.1 + (i // 3) * 2.4),
            Inches(2.8),
            Inches(0.4)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = comp_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Box items
        items_box = slide.shapes.add_textbox(
            Inches(0.4 + (i % 3) * 3.2),
            Inches(start_y + 0.5 + (i // 3) * 2.4),
            Inches(2.8),
            Inches(1.6)
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(comp_items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = RgbColor(255, 255, 255)
    
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ============ SLIDE 1: Title ============
    add_title_slide(
        prs,
        "🎙️ ASIS",
        "Asistent Personal AI Vocal\nAplicație Mobilă Flutter"
    )
    
    # ============ SLIDE 2: Problema și Soluția ============
    add_content_slide(prs, "📌 Problema Adresată", [
        "Utilizatorii pierd timp gestionând task-uri, emails, cumpărături manual",
        "Aplicațiile actuale necesită navigare complexă prin meniuri",
        "Barriere de accesibilitate pentru persoane cu dizabilități sau mâini ocupate",
        "Lipsa unui asistent vocal românesc inteligent și contextual",
        "",
        "## 💡 Soluția: ASIS",
        "- Asistent vocal complet în limba română",
        "- Control hands-free pentru toate funcționalitățile",
        "- AI conversațional cu Google Gemini",
        "- Totul rulează local pe dispozitiv - fără server extern"
    ])
    
    # ============ SLIDE 3: Competiție și Diferențiator ============
    add_content_slide(prs, "🏆 Competiție Globală", [
        "## Asistenți Vocali Integrați",
        "- Google Assistant, Siri, Bixby, Cortana - suport român limitat/inexistent",
        "",
        "## Aplicații AI Third-Party",
        "- ChatGPT (OpenAI) - conversațional, dar fără acțiuni native pe device",
        "- Alexa Mobile - ecosistem Amazon, română nesupinută",
        "- Replika, Character.AI - focus pe companion, nu productivitate",
        "- Rabbit R1, Humane AI Pin - hardware dedicat, preț ridicat",
        "",
        "## Asistenți de Productivitate",
        "- Todoist, Any.do - task management fără voce avansată",
        "- Motion, Reclaim.ai - calendar AI, fără limba română"
    ])
    
    # ============ SLIDE 3b: Diferențiator ============
    add_content_slide(prs, "⭐ Ce Ne Diferențiază", [
        "## Avantaje Unice ASIS",
        "- 100% optimizat pentru limba română nativă",
        "- Arhitectură complet locală (fără server extern, confidențialitate)",
        "- Integrare nativă cu email, calendar, taskuri, cumpărături",
        "- Context conversațional persistent între sesiuni",
        "",
        "## Flexibilitate",
        "- Open source și extensibil prin plugin-uri",
        "- Personalizare completă pentru fiecare utilizator",
        "- Arhitectură modulară pentru adăugare funcționalități noi",
        "",
        "## Viziune pe Termen Lung",
        "- Platformă extensibilă pentru integrări viitoare",
        "- Focus pe productivitate și accesibilitate"
    ])
    
    # ============ SLIDE 4: Arhitectura Tehnică ============
    add_architecture_slide(prs, "🏗️ Arhitectură Tehnică - Flutter/Dart", [
        ("🎤 Voice Services", [
            "SpeechToText Service",
            "TextToSpeech Service",
            "Limba română (ro-RO)"
        ]),
        ("🤖 AI Engine", [
            "Google Gemini 2.5 Flash",
            "Intent Detection",
            "Răspunsuri contextuale"
        ]),
        ("⚡ Action Executor", [
            "Task Management",
            "Shopping Lists",
            "Calendar Events"
        ]),
        ("📧 Email Service", [
            "SMTP/IMAP",
            "Trimitere & Citire",
            "Căutare emailuri"
        ]),
        ("🔍 Search Service", [
            "DuckDuckGo API",
            "Informații real-time",
            "Formatare pentru AI"
        ]),
        ("💾 Local Database", [
            "Hive Database",
            "Offline First",
            "Sincronizare rapidă"
        ])
    ])
    
    # ============ SLIDE 5: Funcționalități ============
    add_content_slide(prs, "✨ Funcționalități Principale", [
        "## Comenzi Vocale & Acțiuni",
        "- \"Adaugă lapte pe lista de cumpărături\"",
        "- \"Creează task: să sun la doctor mâine\"",
        "- \"Trimite email lui Ion cu subiect Întâlnire\"",
        "",
        "## Gestionare Inteligentă",
        "- Task-uri cu priorități și deadline-uri",
        "- Liste de cumpărături pe categorii",
        "- Evenimente calendar cu Google Meet",
        "",
        "## AI Conversațional",
        "- Răspunsuri naturale în română",
        "- Căutare informații pe internet în timp real",
        "- Context conversațional persistent"
    ])
    

    
    # ============ SLIDE 6: Tehnologii ============
    add_content_slide(prs, "🛠️ Stack Tehnologic", [
        "## Frontend & Core",
        "- Flutter 3.x cu Dart",
        "- Material Design 3",
        "- Animații fluide pentru UI/UX",
        "",
        "## Servicii AI & Voice",
        "- google_generative_ai - Gemini API",
        "- speech_to_text - Recunoaștere vocală",
        "- flutter_tts - Sinteză vocală",
        "",
        "## Date & Networking",
        "- Hive - Bază de date NoSQL locală",
        "- mailer - SMTP pentru emails",
        "- http - Căutări internet"
    ])
    
    # ============ SLIDE 7: Demo Flow ============
    add_content_slide(prs, "🎬 Demo - Flow Aplicație", [
        "1️⃣  Utilizatorul apasă butonul microfonului",
        "2️⃣  Speech-to-Text convertește vocea în text",
        "3️⃣  Gemini AI analizează intenția și extrage date",
        "4️⃣  Action Executor execută acțiunea detectată",
        "5️⃣  AI generează răspuns natural personalizat",
        "6️⃣  Text-to-Speech redă răspunsul vocal",
        "",
        "## Exemplu Live",
        "- \"Adaugă 2 kg de mere și pâine pe listă\"",
        "- \"Care sunt taskurile mele pentru azi?\"",
        "- \"Programează o întâlnire cu Ana la 15:00\""
    ])
    
    # ============ SLIDE 8: Rezultate ============
    add_content_slide(prs, "📊 Rezultate Obținute", [
        "## Performanță",
        "- Timp răspuns: < 2 secunde",
        "- Acuratețe recunoaștere vocală: 95%+",
        "- Zero latență server (totul local)",
        "",
        "## Funcțional",
        "- 15+ tipuri de acțiuni suportate",
        "- Suport complet limba română",
        "- Mod offline pentru funcții locale",
        "",
        "## UX",
        "- Interfață intuitivă și minimalistă",
        "- Dark mode suportat",
        "- Animații responsive"
    ])
    
    # ============ SLIDE 9: Monetizare ============
    add_content_slide(prs, "💰 Monetizare și Exploatare", [
        "## Model Freemium",
        "- Versiune gratuită cu funcții de bază",
        "- Premium: integrări avansate, comenzi nelimitate",
        "",
        "## Oportunități B2B",
        "- Licențiere pentru companii (asistent intern)",
        "- White-label pentru aplicații terțe",
        "- API as a Service pentru dezvoltatori",
        "",
        "## Piață Țintă",
        "- 19+ milioane vorbitori nativi de română",
        "- Piață asistenti vocali: $15.8B global (2026)",
        "- Nișă neexploatată pentru limba română"
    ])
    
    # ============ SLIDE 10: Concluzii ============
    add_content_slide(prs, "🚀 Concluzii și Dezvoltare Ulterioară", [
        "## Ce am realizat",
        "- Asistent AI vocal complet funcțional în română",
        "- Arhitectură scalabilă și modulară în Flutter",
        "- Integrare seamless cu email, calendar, taskuri",
        "",
        "## Planuri de dezvoltare - Programări Online",
        "- Programare vocală la clinici, cabinete medicale, dentare",
        "- Completare automată formulare de booking online",
        "- Integrare cu: Doctolib, DOC.ro, Clinica.ro, cabinete custom",
        "- Extindere: restaurante, saloane, service auto",
        "",
        "## Alte Funcționalități Viitoare",
        "- Smart Home (Google Home, Philips Hue)",
        "- Mod offline cu AI local (Llama, Gemma)",
        "- Widget-uri home screen",
        "",
        "## 🎯 Viziune: Cel mai bun asistent vocal românesc!"
    ])
    
    # ============ SLIDE 11: Thank You ============
    add_title_slide(
        prs,
        "Mulțumim!",
        "Întrebări?\n\n🎙️ ASIS - Asistentul Tău Personal"
    )
    
    # Save
    output_path = r"c:\Users\robij\OneDrive\Desktop\Game Changer\ASIS_Prezentare_Concurs.pptx"
    prs.save(output_path)
    print(f"✅ Prezentare salvată: {output_path}")

if __name__ == "__main__":
    create_presentation()
